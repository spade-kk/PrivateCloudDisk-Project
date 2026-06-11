"""
PowerPoint 文件抽取器

支持: .pptx, .ppt
"""
from pathlib import Path

from pptx import Presentation

from core.extractors.base import BaseExtractor, ExtractionResult, split_text, TextChunk
from core.config import settings


class PptxExtractor(BaseExtractor):
    name = "pptx_extractor"

    async def extract(self, file_id: str, path: Path) -> ExtractionResult:
        warnings: list[str] = []
        chunks: list[TextChunk] = []
        all_parts: list[str] = []

        try:
            prs = Presentation(str(path))
            slides = min(len(prs.slides), 300)  # 最多 300 页

            for slide_num, slide in enumerate(prs.slides, start=1):
                slide_texts: list[str] = []

                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            p_text = paragraph.text.strip()
                            if p_text:
                                slide_texts.append(p_text)

                    if shape.has_table:
                        table = shape.table
                        for row in table.rows:
                            row_text = " | ".join(
                                cell.text.strip() for cell in row.cells if cell.text.strip()
                            )
                            if row_text:
                                slide_texts.append(row_text)

                if slide_texts:
                    slide_text = f"Slide {slide_num}: " + " ".join(slide_texts)
                    all_parts.append(slide_text)
                    chunks.extend(
                        split_text(
                            text=slide_text,
                            source_type="pptx",
                            chunk_size=settings.chunk_size_chars,
                            overlap=settings.chunk_overlap_chars,
                            slide=slide_num,
                        )
                    )

            if len(prs.slides) > 300:
                warnings.append(f"PPT 共 {len(prs.slides)} 页，仅处理前 300 页")
        except Exception as e:
            return ExtractionResult(
                file_id=file_id,
                extractor=self.name,
                warnings=[f"PPT 解析失败: {e}"],
            )

        text = "\n".join(all_parts)[:settings.content_max_chars]

        return ExtractionResult(
            file_id=file_id,
            extractor=self.name,
            extractor_version=self.version,
            content_text=text,
            chunks=chunks,
            tags=["presentation", "slides"],
            warnings=warnings,
        )